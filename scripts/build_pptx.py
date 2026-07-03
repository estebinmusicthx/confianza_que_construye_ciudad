# -*- coding: utf-8 -*-
"""
build_pptx.py
Genera `pitch.pptx`: la presentacion del pitch en PowerPoint (16:9).
Arranca con una diapositiva roja con la palabra CONFIANZA en grande y luego
las 8 diapositivas del pitch. Ultima hoja: equipo, integrantes e iconos.

Uso:  python scripts/build_pptx.py
"""
import base64
import io
import os
import re
import qrcode
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BACKUP = os.path.join(ROOT, "index_dashboard_original.html")
OUT = os.path.join(ROOT, "pitch.pptx")
ASSETS = os.path.join(ROOT, "outputs", "_pptx_assets")
os.makedirs(ASSETS, exist_ok=True)
SITE_URL = "https://estebinmusicthx.github.io/confianza_que_construye_ciudad/"

RED = RGBColor(0xE4, 0x00, 0x2B)
DRED = RGBColor(0xC1, 0x00, 0x22)
WINE = RGBColor(0x8A, 0x00, 0x16)
AMBER = RGBColor(0xFF, 0xCC, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x16, 0x16, 0x16)
PINK = RGBColor(0xFF, 0xD3, 0xDC)
MUT = RGBColor(0xC2, 0xBC, 0xB4)
TEAL = RGBColor(0x0F, 0x9D, 0x77)
BLUE = RGBColor(0x1E, 0x5F, 0xA5)
AMBER2 = RGBColor(0xF2, 0xA1, 0x00)
DARKAMBER = RGBColor(0x3A, 0x25, 0x00)
FONT = "Arial"

# ---------------------------------------------------------------- assets
raw = re.findall(r'data:image/[a-zA-Z]+;base64,([A-Za-z0-9+/=]+)',
                 open(BACKUP, encoding="utf-8").read())


def _img(i):
    return Image.open(io.BytesIO(base64.b64decode(raw[i]))).convert("RGBA")


def knockout_white(im):
    px = im.load()
    for y in range(im.height):
        for x in range(im.width):
            r, g, b, a = px[x, y]
            if a > 10:
                px[x, y] = (255, 255, 255, a)
    return im


def strip_gray(im):
    px = im.load()
    for y in range(im.height):
        for x in range(im.width):
            r, g, b, a = px[x, y]
            mx, mn = max(r, g, b), min(r, g, b)
            if a > 0 and 38 <= mn and mx <= 102 and (mx - mn) <= 26:
                px[x, y] = (r, g, b, 0)
    return im


def save_png(im, name):
    p = os.path.join(ASSETS, name)
    im.save(p, "PNG")
    return p


P_DATOS = save_png(knockout_white(_img(0)), "datos.png")
P_BOGOTA = save_png(strip_gray(_img(1)), "bogota.png")
P_DATAJAM = save_png(strip_gray(_img(2)), "datajam.png")
P_INST = save_png(strip_gray(_img(3)), "inst.png")

qr = qrcode.QRCode(border=1, box_size=14)
qr.add_data(SITE_URL)
qr.make(fit=True)
P_QR = save_png(qr.make_image(fill_color="#161616", back_color="white").convert("RGBA"), "qr.png")

# ---------------------------------------------------------------- helpers
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=RED):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def para(tf, first=False, align=PP_ALIGN.LEFT, space_after=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    return p


def run(p, text, size, color, bold=True, italic=False, name=FONT):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = name
    r.font.color.rgb = color
    return r


def card(s, l, t, w, h, fill, radius=0.14):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.16)
    tf.margin_bottom = Inches(0.16)
    return tf


def pic(s, path, l, t, h):
    return s.shapes.add_picture(path, Inches(l), Inches(t), height=Inches(h))


# =============================================================== 1. CONFIANZA
s = slide(RED)
tf = box(s, 0.6, 0.7, 12.13, 0.5)
run(para(tf, True), "DATAJAM 2  ·  DATOS PARA LA TRANSPARENCIA", 15, PINK)
tf = box(s, 0.6, 2.5, 12.13, 2.4, MSO_ANCHOR.MIDDLE)
run(para(tf, True, PP_ALIGN.CENTER), "CONFIANZA", 132, WHITE)
tf = box(s, 0.6, 5.0, 12.13, 0.7)
run(para(tf, True, PP_ALIGN.CENTER), "que construye ciudad", 30, AMBER)

# =============================================================== 2. Portada
s = slide(RED)
pic(s, P_DATOS, 0.85, 0.55, 0.75)
b = pic(s, P_BOGOTA, 0, 0.62, 0.5)
b.left = Inches(13.333 - 0.85 - b.width.inches)
tf = box(s, 0.85, 2.05, 11, 0.45)
run(para(tf, True), "DATAJAM 2  ·  DATOS PARA LA TRANSPARENCIA", 15, AMBER)
tf = box(s, 0.85, 2.55, 10.6, 2.9)
p = para(tf, True)
run(p, "Una ciudad no se construye solo con impuestos. ", 43, WHITE)
run(p, "Se construye con confianza.", 43, AMBER)
tf = box(s, 0.85, 5.75, 11.5, 0.6)
run(para(tf, True), "Equipo SOLUTIONSCITY  ·  Caja de Vivienda Popular (CVP)", 20, PINK, bold=False)

# =============================================================== 3. Problema
s = slide(RED)
tf = box(s, 0.85, 1.15, 11, 0.45)
run(para(tf, True), "EL PROBLEMA", 16, AMBER)
tf = box(s, 0.85, 1.75, 11.3, 1.7)
run(para(tf, True), "La información existe. Pero está rota en pedazos.", 40, WHITE)
tf = box(s, 0.85, 3.9, 10.8, 2.4)
run(para(tf, True), "Aporte voluntario, cumplimiento predial, condiciones sociales y "
    "proyectos ciudadanos viven separados entre entidades y cortes de tiempo. "
    "El ciudadano paga… pero no ve. Y lo que no se ve, no construye confianza.", 21, WHITE, bold=False)

# =============================================================== 4. El giro
s = slide(RED)
tf = box(s, 0.85, 2.2, 11.3, 0.7)
run(para(tf, True), "Dejamos de preguntar cuánto más puede aportar la gente.", 21, PINK, bold=False)
tf = box(s, 0.85, 2.95, 11.8, 2.6)
p = para(tf, True)
run(p, "¿Cuánto más puede ", 56, WHITE)
run(p, "explicar", 56, AMBER)
run(p, " la ciudad?", 56, WHITE)

# =============================================================== 5. ITCC
s = slide(RED)
tf = box(s, 0.85, 1.1, 11, 0.45)
run(para(tf, True), "LA INNOVACIÓN", 16, AMBER)
tf = box(s, 0.85, 1.6, 11.4, 1.5)
run(para(tf, True), "ITCC: un mapa de confianza, no un ranking de culpables.", 34, WHITE)
tf = box(s, 0.85, 3.15, 11.4, 1.5)
run(para(tf, True), "Mide, de 0 a 100 y con las mismas reglas para las 20 localidades, qué tan "
    "favorable es cada territorio para hablar de confianza. No juzga personas. "
    "No inventa causas. Mide condiciones.", 19, WHITE, bold=False)
stats = [("0–100", "escala comparable"), ("20", "localidades, mismas reglas"), ("4", "fuentes públicas")]
for i, (n, lab) in enumerate(stats):
    tf = card(s, 0.85 + i * 4.0, 5.05, 3.7, 1.55, DRED)
    run(para(tf, True), n, 40, WHITE)
    run(para(tf), lab, 14, PINK, bold=False)

# =============================================================== 6. 4 Bogotás
s = slide(RED)
tf = box(s, 0.85, 0.95, 11, 0.45)
run(para(tf, True), "EL MODELO DE DECISIÓN", 16, AMBER)
tf = box(s, 0.85, 1.45, 11.4, 1.0)
run(para(tf, True), "No hay una Bogotá. Hay cuatro.", 38, WHITE)
quads = [
    ("Embajadores", "Alta confianza y alto cumplimiento: activar vocerías y reconocimiento.", WINE, WHITE, PINK),
    ("Comunicar", "Buenas condiciones, menor cumplimiento: mostrar resultados y beneficios.", AMBER2, DARKAMBER, DARKAMBER),
    ("Reconocer", "Cumplen, pero la confianza es frágil: agradecer y reforzar transparencia.", TEAL, WHITE, WHITE),
    ("Transparencia", "Más fricción: pedagogía, trazabilidad y presencia institucional.", BLUE, WHITE, WHITE),
]
pos = [(0.85, 2.75), (7.0, 2.75), (0.85, 4.95), (7.0, 4.95)]
for (title, desc, fill, tc, dc), (l, t) in zip(quads, pos):
    tf = card(s, l, t, 5.48, 2.0, fill)
    run(para(tf, True, space_after=4), title, 24, tc)
    run(para(tf), desc, 16, dc, bold=False)

# =============================================================== 7. Propuesta
s = slide(RED)
tf = box(s, 0.85, 1.05, 11, 0.45)
run(para(tf, True), "LA PROPUESTA", 16, AMBER)
tf = box(s, 0.85, 1.55, 11.6, 1.1)
p = para(tf, True)
run(p, "Programa ", 34, WHITE)
run(p, "Confianza que Construye Ciudad", 34, AMBER)
lines = [
    ("1", "Trazabilidad ciudadana", "Mostrar abiertamente cómo se usa el aporte."),
    ("2", "Reconocimiento inteligente", "Del gesto (la tarjeta del Metro) a política permanente."),
    ("3", "Pedagogía diferenciada", "Un mensaje por territorio, no una campaña genérica."),
    ("4", "ITCC vivo", "Un índice que se actualiza con datos abiertos."),
    ("5", "Gobierno abierto", "Decidir con evidencia y límites declarados."),
]
cw, gap = 2.24, 0.19
start = (13.333 - (5 * cw + 4 * gap)) / 2
for i, (n, title, desc) in enumerate(lines):
    l = start + i * (cw + gap)
    tf = card(s, l, 3.3, cw, 2.9, DRED)
    run(para(tf, True, space_after=6), n, 30, AMBER)
    run(para(tf, space_after=4), title, 16, WHITE)
    run(para(tf), desc, 12.5, PINK, bold=False)

# =============================================================== 8. Cierre
s = slide(RED)
tf = box(s, 0.85, 1.5, 8.0, 2.4)
p = para(tf, True)
run(p, "La confianza no se decreta. ", 44, WHITE)
run(p, "Se demuestra.", 44, AMBER)
cl = [("185.485", "contribuyentes"), ("$7,1 mil M", "aportados (COP)"), ("20", "localidades")]
for i, (n, lab) in enumerate(cl):
    tf = card(s, 0.85 + i * 2.65, 4.35, 2.45, 1.35, DRED)
    run(para(tf, True), n, 30, WHITE)
    run(para(tf), lab, 13, PINK, bold=False)
qcard = card(s, 9.75, 2.15, 2.75, 2.75, WHITE)
pic(s, P_QR, 10.02, 2.42, 2.2)
tf = box(s, 9.75, 5.0, 2.75, 0.4)
run(para(tf, True, PP_ALIGN.CENTER), "Explórelo →", 15, WHITE)

# =============================================================== 9. Creditos
s = slide(INK)
tf = box(s, 0.85, 0.85, 11, 0.45)
run(para(tf, True), "EQUIPO", 16, AMBER)
tf = box(s, 0.85, 1.35, 11.6, 1.1)
run(para(tf, True), "SOLUTIONSCITY", 48, WHITE)
tf = box(s, 0.85, 2.5, 11.6, 0.5)
run(para(tf, True), "en representación de la Caja de Vivienda Popular — CVP", 22, PINK, bold=False)
members = [
    ("Julio Esteban Fuentes Herrera", "  —  Perfil técnico en análisis y visualización de datos"),
    ("Juan Carlos Sanabria Medina", "  —  Perfil de análisis sectorial o de política pública"),
    ("Diego Armando Gonzalez Quiroga", "  —  Perfil complementario (temático, metodológico o técnico)"),
]
tf = box(s, 0.85, 3.5, 11.8, 2.0)
for i, (nom, rol) in enumerate(members):
    p = para(tf, i == 0, space_after=10)
    run(p, nom, 24, WHITE)
    run(p, rol, 16, MUT, bold=False)
pic(s, P_DATAJAM, 0.85, 6.15, 0.95)
pic(s, P_DATOS, 3.1, 6.25, 0.7)
pic(s, P_BOGOTA, 4.5, 6.35, 0.5)
pic(s, P_INST, 6.2, 6.4, 0.62)

prs.save(OUT)
print("OK -> pitch.pptx  (%d diapositivas, %d KB)" % (len(prs.slides._sldIdLst), os.path.getsize(OUT) // 1024))
